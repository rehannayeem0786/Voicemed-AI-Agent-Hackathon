#!/usr/bin/env python3
"""Generate the VoiceMed AI hackathon pitch deck.

    python scripts/make_deck.py   ->  VoiceMed_AI_Pitch_Deck.pptx

12 slides, 16:9, styled to match the app's clinical navy + teal theme.
Requires: pip install python-pptx
"""

import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")

# ─── Theme (mirrors static/style.css — v2 dark-glass redesign) ──
BG = RGBColor.from_string("070B14")
PANEL = RGBColor.from_string("111A2D")
PANEL2 = RGBColor.from_string("18233A")
LINE = RGBColor.from_string("24334C")
TEXT = RGBColor.from_string("E9EEF7")
MUTED = RGBColor.from_string("8EA1BB")
TEAL = RGBColor.from_string("2DD4BF")
SKY = RGBColor.from_string("38BDF8")
RED = RGBColor.from_string("F43F5E")
AMBER = RGBColor.from_string("F59E0B")
GREEN = RGBColor.from_string("34D399")
INK = RGBColor.from_string("04211C")  # dark text on accent fills

FONT = "Segoe UI"
FONT_SB = "Segoe UI Semibold"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
prs.core_properties.title = "VoiceMed AI — Voice-Powered Medical Triage Nurse"
prs.core_properties.author = "Rehan Nayeem"
BLANK = prs.slide_layouts[6]


# ─── Helpers ────────────────────────────────────────────────────
def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def box(s, x, y, w, h, fill=PANEL, line=LINE, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10):
    b = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            b.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        b.fill.background()
    else:
        b.fill.solid()
        b.fill.fore_color.rgb = fill
    if line is None:
        b.line.fill.background()
    else:
        b.line.color.rgb = line
        b.line.width = Pt(lw)
    b.shadow.inherit = False
    return b


def text(s, x, y, w, h, content, size=14, color=TEXT, bold=False, italic=False, font=FONT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, wrap=True):
    """content: str | list of paragraphs; a paragraph is str | (str, overrides) |
    list of (str, overrides) run-tuples."""
    t = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = t.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paragraphs = content if isinstance(content, list) else [content]
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        runs = para if isinstance(para, list) else [para if isinstance(para, tuple) else (para, {})]
        for rtext, ov in runs:
            r = p.add_run()
            r.text = rtext
            f = r.font
            f.name = ov.get("font", font)
            f.size = Pt(ov.get("size", size))
            f.bold = ov.get("bold", bold)
            f.italic = ov.get("italic", italic)
            f.color.rgb = ov.get("color", color)
    return t


def kicker(s, label, sub=None):
    text(s, 0.62, 0.42, 11.5, 0.32, label, size=12.5, color=TEAL, bold=True, font=FONT_SB)
    if sub:
        text(s, 0.62, 0.76, 12.1, 0.85, sub, size=30, bold=True, font=FONT_SB)


def chip(s, x, y, w, h, label, fill=TEAL, color=INK, size=11.5, bold=True, line=None):
    c = box(s, x, y, w, h, fill=fill, line=line, radius=0.5)
    text(s, x, y, w, h, label, size=size, color=color, bold=bold,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return c


def dot_bullets(s, x, y, w, items, gap=0.52, size=13.5, dot=TEAL):
    for i, item in enumerate(items):
        yy = y + i * gap
        box(s, x, yy + 0.085, 0.11, 0.11, fill=dot, line=None, shape=MSO_SHAPE.OVAL)
        if isinstance(item, tuple):
            head, body = item
            text(s, x + 0.28, yy, w - 0.28, gap,
                 [(head + "  ", {"bold": True, "size": size}), (body, {"color": MUTED, "size": size - 0.5})],
                 size=size)
        else:
            text(s, x + 0.28, yy, w - 0.28, gap, item, size=size)


def waveform(s, x0, x1, ycenter, n=90, amp=0.42, base=0.05):
    import math
    step = (x1 - x0) / n
    for i in range(n):
        h = base + abs(math.sin(i * 0.55)) * amp * (0.35 + 0.65 * abs(math.sin(i * 0.11)))
        box(s, x0 + i * step, ycenter - h / 2, max(0.014, step * 0.45), h,
            fill=TEAL if i % 3 else SKY, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)


def arrow(s, x, y, w, h, color=LINE, flip=False):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    if flip:
        a.rotation = 180
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt


# ─── Slide 1 · Hero ─────────────────────────────────────────────
s = slide()
waveform(s, 0.6, 12.75, 6.95, n=110, amp=0.55)
box(s, 0.6, 6.55, 12.13, 0.02, fill=LINE, line=None, shape=MSO_SHAPE.RECTANGLE)
chip(s, 0.62, 0.62, 4.55, 0.44, "🏆  ASSEMBLYAI VOICE AGENT HACKATHON", fill=PANEL2, color=TEAL, size=12)
text(s, 0.6, 1.7, 12.1, 1.3,
     [[("VoiceMed ", {"color": TEXT}), ("AI", {"color": TEAL})]], size=63, bold=True, font=FONT_SB)
text(s, 0.62, 3.0, 11.6, 0.6,
     [[("Aria — ", {"bold": True}), ("the AI triage nurse you talk to, not type at.", {"color": MUTED})]], size=22)
text(s, 0.62, 3.85, 11.9, 1.15,
     [[("Voice intake → ESI 1–5 triage → drug-interaction checks → emergency escalation → EMR-ready SOAP notes.", {})],
      [("One real-time voice connection, in English & Español — built on the AssemblyAI Voice Agent API.",
        {"color": MUTED})]], size=14.5, spacing=1.3)
text(s, 0.62, 7.0, 12.1, 0.4,
     [[("github.com/rehannayeem0786/Voicemed-AI-Agent-Hackathon   ·   ", {"color": MUTED, "size": 11}),
       ("protocol verified live against wss://agents.assemblyai.com", {"color": GREEN, "size": 11, "bold": True})]])
notes(s, "Open: 'Hi judges — think back to the last time you tried to describe pain to a health system. Forms, "
         "hold music, apps. VoiceMed replaces that with a conversation.' Everything on screen runs on ONE "
         "AssemblyAI voice WebSocket, and the protocol was verified live against the real endpoint.")

# ─── Slide 2 · Problem ──────────────────────────────────────────
s = slide()
kicker(s, "01 · THE PROBLEM", "Every triage call starts with a keyboard")
stats = [("84%", "of adults 60+ struggle to use online patient portals"),
         ("60%", "of US adults lack the health-literacy skills to navigate forms"),
         ("4.5B", "GP / telehealth visits a year worldwide still begin as text intake")]
for i, (num, lbl) in enumerate(stats):
    x = 0.62 + i * 4.08
    box(s, x, 1.85, 3.86, 1.72)
    text(s, x + 0.3, 2.06, 3.3, 0.8, num, size=40, bold=True, color=TEAL, font=FONT_SB)
    text(s, x + 0.3, 2.86, 3.3, 0.65, lbl, size=12, color=MUTED, spacing=1.05)
box(s, 0.62, 3.95, 12.11, 2.6, fill=PANEL2)
text(s, 0.95, 4.2, 11.5, 0.4, "What breaks in a real emergency", size=15, bold=True, color=AMBER)
dot_bullets(s, 0.95, 4.72, 11.3, [
    ("Typing ≠ triage.", "Red-flag symptoms get flattened into dropdowns — chest pain becomes 'general discomfort'."),
    ("Language barriers.", "56M+ Spanish speakers navigate English-only intake at the moment of highest stress."),
    ("The 3 a.m. gap.", "When someone is dizzy, sweating and scared, the safest interface is a calm voice asking one question at a time."),
], gap=0.58, size=13.5)
text(s, 0.62, 6.85, 12.1, 0.4,
     [[("Voice is the only UI that works ", {"color": MUTED, "italic": True}),
       ("at zero mobility, zero literacy, zero spare attention.", {"color": TEAL, "italic": True, "bold": True})]])
notes(s, "Land three numbers, then the punchline: in an emergency the user has zero mobility, zero literacy "
         "and zero spare attention — only voice fits. This motivates a VOICE-FIRST agent, not another chatbot.")

# ─── Slide 3 · Live demo ────────────────────────────────────────
s = slide()
kicker(s, "02 · LIVE DEMO", "Ninety seconds, one take — the Demo Copilot script")
beats = [
    ("🚨", "Chest pain", "\u201cCrushing chest pain into my left arm, I'm sweating.\u201d",
     "red-flag fired \u2192 ESI 1 \u2192 Aria says call 911", RED),
    ("💊", "Medications", "\u201cI also take warfarin and aspirin.\u201d",
     "interaction check \u2192 high bleeding risk panel", AMBER),
    ("🔍", "Severity", "\u201cStarted this morning, it's 8 out of 10.\u201d",
     "symptom lookup \u2192 follow-ups woven into Aria's questions", SKY),
    ("📋", "Wrap-up", "\u201cI think that's everything.\u201d",
     "ESI gauge animates \u2192 SOAP note with ICD-10 \u2192 booking", GREEN),
]
for i, (icon, ttl, say, watch, accent) in enumerate(beats):
    x = 0.62 + i * 3.08
    box(s, x, 1.95, 2.86, 3.9)
    box(s, x, 1.95, 2.86, 0.14, fill=accent, line=None, shape=MSO_SHAPE.RECTANGLE)
    text(s, x + 0.24, 2.28, 2.4, 0.5, icon + "  " + ttl, size=15, bold=True, font=FONT_SB)
    text(s, x + 0.24, 2.86, 2.4, 1.05, say, size=12, color=TEXT, italic=True, spacing=1.1)
    text(s, x + 0.24, 4.05, 2.4, 0.3, "WATCH", size=10, color=MUTED, bold=True)
    text(s, x + 0.24, 4.35, 2.4, 1.3, watch, size=11.5, color=accent, spacing=1.1)
    if i < 3:
        arrow(s, x + 2.86, 3.6, 0.24, 0.32, color=LINE)
box(s, 0.62, 6.2, 12.11, 0.75, fill=PANEL2)
text(s, 0.92, 6.33, 11.6, 0.5,
     [[("🎬 Demo Copilot  ", {"color": TEAL, "bold": True, "size": 13}),
       ("built into the app (Alt+D): auto-advances as Aria hears each line — a flawless one-take demo, even in a loud room.",
        {"color": MUTED, "size": 12.5})]])
notes(s, "If doing the live demo, switch to browser now and press Alt+D. If recording, show the 4 beats in order: "
         "emergency banner, medication panel, ESI gauge + SOAP note. The Copilot card on the left mirrors this slide.")

# ─── Slide 4 · Architecture ─────────────────────────────────────
s = slide()
kicker(s, "03 · ARCHITECTURE", "One WebSocket in — everything else is local")
box(s, 0.62, 1.8, 3.4, 4.35, fill=PANEL2)
text(s, 0.86, 2.0, 3.0, 0.4, "Browser", size=15, bold=True, color=SKY, font=FONT_SB)
dot_bullets(s, 0.86, 2.5, 2.95, [
    "Mic capture, VAD, barge-in",
    "Animated waveform + live transcript",
    "Progress tracker & tool panels",
    " PCM 16 kHz down · 24 kHz up",
], gap=0.62, size=12.5, dot=SKY)
chip(s, 0.9, 5.35, 2.8, 0.5, "FastAPI bridge", fill=PANEL, color=TEXT, size=12.5, line=LINE)
arrow(s, 4.14, 3.75, 0.5, 0.42, color=TEAL)
box(s, 4.72, 1.8, 3.85, 4.35, fill=PANEL2)
box(s, 4.72, 1.8, 3.85, 0.12, fill=TEAL, line=None, shape=MSO_SHAPE.RECTANGLE)
text(s, 4.96, 2.05, 3.4, 0.45, "AssemblyAI Voice Agent API", size=15, bold=True, color=TEAL, font=FONT_SB)
text(s, 4.96, 2.5, 3.4, 0.3, "wss://agents.assemblyai.com", size=10.5, color=MUTED)
dot_bullets(s, 4.96, 2.95, 3.4, [
    ("Universal STT", "streaming, EN + ES"),
    ("LLM tool-calling", "session.update contract"),
    ("Dialog TTS", "voice alba / lola"),
], gap=0.6, size=12.5)
chip(s, 5.0, 5.0, 3.3, 0.5, "auth: short-lived token API", fill=BG, color=MUTED, size=10.5, line=LINE)
chip(s, 5.0, 5.6, 3.3, 0.5, "or stored agent_id binding", fill=BG, color=MUTED, size=10.5, line=LINE)
arrow(s, 8.69, 3.75, 0.5, 0.42, color=TEAL)
box(s, 9.31, 1.8, 3.42, 4.35, fill=PANEL2)
text(s, 9.55, 2.0, 3.0, 0.4, "Local clinical stack", size=15, bold=True, color=GREEN, font=FONT_SB)
dot_bullets(s, 9.55, 2.5, 2.95, [
    "SQLite session + audit log",
    "JSON clinical databases",
    "6 function tools in Python",
    "SOAP notes, JSON export",
], gap=0.62, size=12.5, dot=GREEN)
text(s, 0.62, 6.45, 12.1, 0.6,
     [[("Why this shape: ", {"bold": True, "color": TEXT}),
       ("client-side tools respond in milliseconds (no LLM round-trip for lookups), PHI never leaves the bridge, and the agent works with or without a published stored agent.",
        {"color": MUTED})]], size=12.5, spacing=1.15)
notes(s, "Walk left to right: audio in the browser, one websocket to AssemblyAI for STT+LLM+TTS, tool calls "
         "executed locally against our clinical data, results queued back as tool.result. Emphasize latency + privacy.")

# ─── Slide 5 · Agent anatomy ────────────────────────────────────
s = slide()
kicker(s, "04 · THE AGENT", "Aria is designed like a real triage nurse")
box(s, 0.62, 1.8, 5.9, 4.5, fill=PANEL2)
text(s, 0.9, 2.02, 5.3, 0.4, "Stored-agent definition (agents/voicemed_triage.jsonc)", size=13, bold=True, color=TEAL)
dot_bullets(s, 0.9, 2.55, 5.35, [
    ("One question at a time.", "Voice rules stop the LLM from lecturing — every turn is short and ends listening."),
    ("Warm, plain-language persona.", "No jargon, no lists, brief empathy before the next question."),
    ("Safety is procedural, not hoped for.", "Red-flag screening is a numbered interview stage, escalation is a held tool call."),
    ("Structured by design.", "Symptom → meds → red flags → ESI → SOAP → booking, mirrored by the UI tracker."),
], gap=0.92, size=13)
box(s, 6.82, 1.8, 5.9, 4.5, fill=PANEL2)
text(s, 7.1, 2.02, 5.3, 0.4, "Why judges should care", size=13, bold=True, color=AMBER)
dot_bullets(s, 7.1, 2.55, 5.35, [
    ("Reproducible agent.", "The JSONC is the single source of truth for the browser, the publisher script and tests."),
    ("Deterministic safety net.", "The ESI scorer and interaction checker are local code — the LLM can talk, but code decides."),
    ("Instrumented UI.", "Every tool call surfaces in panels — nothing happens invisibly."),
    ("Language parity.", "agents/voicemed_triage_es.jsonc: same tools, Spanish prompt, voice lola."),
], gap=0.92, size=13, dot=AMBER)
text(s, 0.62, 6.55, 12.1, 0.5,
     [[("An agent is only as safe as its process — ", {"italic": True, "color": MUTED}),
       ("so we made the process explicit, testable and visible.", {"italic": True, "bold": True, "color": TEAL})]])
notes(s, "Differentiator slide: most hackathon agents are a one-line prompt. Ours is a reviewable, versioned "
         "definition with voice rules, staged screening and held tool calls — plus local deterministic scoring.")

# ─── Slide 6 · Clinical tools ───────────────────────────────────
s = slide()
kicker(s, "05 · CLINICAL TOOLS", "Six function tools — executed locally, results spoken back")
tools = [
    ("symptom_lookup", "Matches symptoms to conditions, red flags & suggested follow-ups", SKY),
    ("drug_interaction_check", "Screens med list against a clinical interaction database", AMBER),
    ("triage_assessment", "Scores ESI 1–5 from severity, duration, vitals & red flags", TEAL),
    ("generate_soap_note", "Emits a structured S/O/A/P note with ICD-10 codes", GREEN),
    ("book_appointment", "Books telehealth / urgent care / specialist follow-up", SKY),
    ("emergency_alert", "Holds the session & escalates 911 / 988 on red flags", RED),
]
for i, (name, desc, accent) in enumerate(tools):
    col, row = i % 3, i // 3
    x, y = 0.62 + col * 4.08, 1.95 + row * 1.95
    box(s, x, y, 3.86, 1.72)
    box(s, x, y, 0.12, 1.72, fill=accent, line=None, shape=MSO_SHAPE.RECTANGLE)
    text(s, x + 0.3, y + 0.2, 3.4, 0.35, name, size=13.5, bold=True, color=accent, font="Consolas")
    text(s, x + 0.3, y + 0.62, 3.35, 1.0, desc, size=11.5, color=MUTED, spacing=1.1)
box(s, 0.62, 5.95, 12.11, 1.0, fill=PANEL2)
text(s, 0.92, 6.1, 11.6, 0.7,
     [[("Tool loop: ", {"bold": True, "color": TEXT}),
       ("agent calls tool → bridge executes in ms → result queued as tool.result (held until reply.done) → Aria speaks the finding in one sentence — never raw JSON.",
        {"color": MUTED})]], size=12.5, spacing=1.2)
notes(s, "Six tools cover the whole intake lifecycle. Highlight emergency_alert: execution_mode hold keeps the "
         "patient with Aria while the alert lands. Tool names stay English so both languages share one dispatcher.")

# ─── Slide 7 · Bilingual ────────────────────────────────────────
s = slide()
kicker(s, "06 · BILINGUAL TRIAGE", "English · Español — same agent, same tools, one click")
box(s, 0.62, 1.8, 5.9, 4.15, fill=PANEL2)
box(s, 0.62, 1.8, 5.9, 0.12, fill=SKY, line=None, shape=MSO_SHAPE.RECTANGLE)
text(s, 0.9, 2.05, 5.3, 0.45, "🇺🇸  English — voice: alba", size=15, bold=True, color=SKY, font=FONT_SB)
text(s, 0.9, 2.6, 5.35, 0.75, "\u201cWhat's the main reason for your call today?\u201d",
     size=13.5, italic=True, color=TEXT)
dot_bullets(s, 0.9, 3.35, 5.35, [
    "English medical keyterms prompt",
    "ESI scoring, red flags, 911 / 988 escalation",
    "SOAP note for the medical record",
], gap=0.62, size=12.5, dot=SKY)
box(s, 6.82, 1.8, 5.9, 4.15, fill=PANEL2)
box(s, 6.82, 1.8, 5.9, 0.12, fill=TEAL, line=None, shape=MSO_SHAPE.RECTANGLE)
text(s, 7.1, 2.05, 5.3, 0.45, "🇪🇸  Español — voice: lola", size=15, bold=True, color=TEAL, font=FONT_SB)
text(s, 7.1, 2.6, 5.35, 0.75, "\u201c¿Cuál es el motivo principal de tu llamada hoy?\u201d",
     size=13.5, italic=True, color=TEXT)
dot_bullets(s, 7.1, 3.35, 5.35, [
    "Spanish STT (language_codes: es) + keyterms",
    "Full Spanish system prompt & greeting",
    "Native Spanish TTS voice — not translated output",
], gap=0.62, size=12.5)
box(s, 0.62, 6.15, 12.11, 0.85, fill=PANEL2)
text(s, 0.92, 6.3, 11.6, 0.55,
     [[("Design choice: ", {"bold": True}),
       ("patients in crisis must not think in a second language — the interview, follow-ups and safety instructions all happen in their language, while tool results stay structured for the clinic.",
        {"color": MUTED})]], size=12.5, spacing=1.15)
notes(s, "Switch the 🌐 selector in the app to demo live. Point out voice lola is a native Spanish voice, "
         "not English TTS reading Spanish — and that the tool layer is shared, so results stay identical.")

# ─── Slide 8 · Safety & trust ───────────────────────────────────
s = slide()
kicker(s, "07 · SAFETY & TRUST", "Medical AI earns trust by failing loudly and safely")
safety = [
    ("🚨", "Red flags first", "Stroke, cardiac, bleeding, anaphylaxis and suicidal-ideation screening is a mandatory interview stage — escalation is a held tool call, not a suggestion.", RED),
    ("🧮", "Code decides severity", "ESI 1–5 comes from a deterministic local scorer; the LLM narrates, it never grades urgency.", TEAL),
    ("🚫", "No diagnosis, ever", "Hard prompt rules: no dosing advice, no prescriptions, 'could be related to' instead of 'you have'.", AMBER),
    ("🔒", "Privacy by design", "Audio streams through the agent bridge; clinical records stay local. No PHI in prompts to third parties.", GREEN),
]
for i, (icon, ttl, body, accent) in enumerate(safety):
    col, row = i % 2, i // 2
    x, y = 0.62 + col * 6.2, 1.9 + row * 2.1
    box(s, x, y, 5.9, 1.88)
    box(s, x, y, 5.9, 0.12, fill=accent, line=None, shape=MSO_SHAPE.RECTANGLE)
    text(s, x + 0.28, y + 0.24, 5.3, 0.4, icon + "  " + ttl, size=15, bold=True, font=FONT_SB, color=accent)
    text(s, x + 0.28, y + 0.72, 5.35, 1.05, body, size=12, color=MUTED, spacing=1.12)
text(s, 0.62, 6.35, 12.1, 0.5,
     [[("Disclaimer, honestly stated: ", {"bold": True, "color": AMBER}),
       ("a demo, not a medical device — but the safety architecture is the part worth grading.", {"color": MUTED})]])
notes(s, "Judges probe safety on medical projects. Answer preemptively: mandatory red-flag stage, deterministic "
         "ESI scoring, no-diagnosis rules, and a held emergency_alert call. End with the honest disclaimer.")

# ─── Slide 9 · Engineering quality ──────────────────────────────
s = slide()
kicker(s, "08 · ENGINEERING", "Not a demo hack — a verified codebase")
eq = [
    ("44", "unit & API tests passing (pytest)"), ("7/7", "live protocol checks vs the real endpoint"),
    ("6", "function tools with typed schemas"), ("2", "languages with full agent parity"),
]
for i, (num, lbl) in enumerate(eq):
    x = 0.62 + i * 3.08
    box(s, x, 1.85, 2.86, 1.5)
    text(s, x + 0.26, 2.0, 2.4, 0.65, num, size=30, bold=True, color=TEAL, font=FONT_SB)
    text(s, x + 0.26, 2.68, 2.4, 0.6, lbl, size=11.5, color=MUTED, spacing=1.05)
dot_bullets(s, 0.62, 3.7, 6.0, [
    ("Smoke test against production.", "connects to wss://agents.assemblyai.com: token → session.update → ready → audio → clean end."),
    ("One config, three consumers.", "the JSONC drives the browser session, the publisher script and the tests."),
    ("Async FastAPI bridge.", "lifespan-managed, env-driven, typed tool dispatch."),
], gap=0.72, size=13)
box(s, 6.92, 3.62, 5.8, 2.5, fill=PANEL2)
text(s, 7.2, 3.82, 5.3, 0.35, "Verified live, in order", size=13, bold=True, color=GREEN)
dot_bullets(s, 7.2, 4.25, 5.3, [
    "token mint (single-use, 2-min TTL)",
    "session.update (agent or inline config)",
    "session.ready → greeting audio bytes",
    "graceful session.end, zero leaks",
], gap=0.45, size=11.5, dot=GREEN)
notes(s, "This is the 'we did our homework' slide: 44 tests, plus a smoke test that replays the actual Voice "
         "Agent protocol against production and prints PASS — judges can run it in one command.")

# ─── Slide 10 · Impact & roadmap ────────────────────────────────
s = slide()
kicker(s, "09 · IMPACT", "Where this goes the morning after the hackathon")
box(s, 0.62, 1.8, 5.9, 4.55, fill=PANEL2)
text(s, 0.9, 2.02, 5.3, 0.4, "Why teams adopt it", size=15, bold=True, color=TEAL, font=FONT_SB)
dot_bullets(s, 0.9, 2.6, 5.35, [
    ("Staff leverage.", "A pre-nurse does structured intake so human nurses start at level 3, not level 0."),
    ("Fewer mis-triaged calls.", "Mandatory red-flag stage + deterministic scoring cut the 'I didn't know it was serious' failure mode."),
    ("Health equity.", "Voice + Spanish reaches the patients portals miss."),
], gap=0.98, size=13)
box(s, 6.82, 1.8, 5.9, 4.55, fill=PANEL2)
text(s, 7.1, 2.02, 5.3, 0.4, "Roadmap (all local, plausible)", size=15, bold=True, color=AMBER, font=FONT_SB)
steps = ["EHR write-back (FHIR) for the SOAP note",
         "Pill-image & vitals capture via device APIs",
         "Voice-keyterm dashboard for clinic load",
         "On-prem model swap behind the same tools"]
for i, st in enumerate(steps):
    yy = 2.62 + i * 0.92
    chip(s, 7.1, yy, 0.62, 0.5, f"{i+1}", fill=PANEL, color=TEAL, size=14, line=TEAL)
    text(s, 7.92, yy + 0.02, 4.6, 0.8, st, size=12.5, spacing=1.05)
notes(s, "Keep the roadmap concrete and unfancy: FHIR write-back, device vitals, clinic dashboard, on-prem "
         "swap. Each maps to an existing seam in the code — tools stay, models change.")

# ─── Slide 11 · What's shipped ──────────────────────────────────
s = slide()
kicker(s, "10 · SHIPPED", "Everything below is in the repo — running today")
shipped = [
    ("🎙️", "Real-time voice agent", "Barge-in, VAD, live transcript, animated waveform"),
    ("🩺", "Full triage flow", "ESI scoring, red flags, med checks, SOAP, booking"),
    ("🌍", "English + Español", "Language switcher, native voices, shared tools"),
    ("🎬", "Demo Copilot", "Presenter mode that auto-advances with the patient"),
    ("📊", "Clinical dashboard", "Redesigned dark-glass UI: chat transcript, stepper, ESI gauge, tool timeline"),
    ("🧪", "44 tests + smoke", "Protocol verified against production endpoints"),
]
for i, (icon, ttl, body) in enumerate(shipped):
    col, row = i % 3, i // 3
    x, y = 0.62 + col * 4.08, 1.9 + row * 2.0
    box(s, x, y, 3.86, 1.78)
    text(s, x + 0.28, y + 0.2, 3.3, 0.42, icon + "  " + ttl, size=14, bold=True, font=FONT_SB, color=TEXT)
    text(s, x + 0.28, y + 0.72, 3.32, 0.95, body, size=11.5, color=MUTED, spacing=1.1)
text(s, 0.62, 6.2, 12.1, 0.6,
     [[("Run it: ", {"bold": True, "color": TEAL}),
       ("pip install -r requirements.txt  →  add ASSEMBLYAI_API_KEY  →  uvicorn app.main:app  →  press Start.",
        {"color": MUTED})]], size=13)
notes(s, "Recap grid — speak over it in 20 seconds. Then the one-liner: clone, key, uvicorn, press Start.")

# ─── Slide 12 · Closing ─────────────────────────────────────────
s = slide()
waveform(s, 0.6, 12.75, 1.15, n=110, amp=0.35)
text(s, 0.6, 1.9, 12.1, 1.1,
     [[("The next health call ", {"color": TEXT}), ("should sound", {"color": TEXT}),
       (" like a conversation.", {"color": TEAL})]], size=36, bold=True, font=FONT_SB)
text(s, 0.62, 3.05, 11.9, 0.6,
     "VoiceMed AI — Aria, the bilingual triage nurse on the AssemblyAI Voice Agent API.",
     size=16, color=MUTED)
box(s, 0.62, 4.0, 12.11, 1.9, fill=PANEL2)
text(s, 0.95, 4.2, 11.4, 0.4, "Try it in 60 seconds", size=13, bold=True, color=TEAL)
text(s, 0.95, 4.62, 11.4, 1.2,
     [[("git clone github.com/rehannayeem0786/Voicemed-AI-Agent-Hackathon", {"font": "Consolas", "size": 13})],
      [(".env  ←  ASSEMBLYAI_API_KEY=...      uvicorn app.main:app --port 8000", {"font": "Consolas", "size": 13, "color": MUTED})],
      [("open http://localhost:8000  →  press Start  →  say \u201cI have crushing chest pain\u201d", {"font": "Consolas", "size": 13, "color": SKY})]],
     size=13, spacing=1.3)
text(s, 0.62, 6.25, 12.1, 0.5,
     [[("Rehan Nayeem", {"bold": True, "size": 13}),
       ("   ·   AssemblyAI Voice Agent Hackathon   ·   thank you — questions welcome 🎤", {"color": MUTED, "size": 13})]])
notes(s, "Close: restate the one-liner, point at the repo QR/URL, thank judges, invite the chest-pain question "
         "live — that demo beat never misses. 30 seconds total.")

# ─── Save ───────────────────────────────────────────────────────
OUT = "VoiceMed_AI_Pitch_Deck.pptx"
prs.save(OUT)
print(f"✅ Saved {OUT} with {len(prs.slides._sldIdLst)} slides")
print("   16:9, speaker notes on every slide")